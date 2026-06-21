from pathlib import Path

from pptx import Presentation

from app.loaders.base import BaseLoader
from app.rag.schema import RawDocument


class PPTLoader(BaseLoader):
    supported_extensions = (".pptx", ".ppt")

    def load(self, file_path: Path) -> list[RawDocument]:
        if file_path.suffix.lower() == ".ppt":
            raise NotImplementedError("Legacy .ppt 文件建议先转换为 .pptx 后再解析。")

        documents: list[RawDocument] = []
        prs = Presentation(file_path)
        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                documents.append(
                    RawDocument(
                        content="\n".join(texts),
                        source_file=file_path.name,
                        file_type="pptx",
                        metadata={"slide": slide_idx},
                    )
                )
        return documents
